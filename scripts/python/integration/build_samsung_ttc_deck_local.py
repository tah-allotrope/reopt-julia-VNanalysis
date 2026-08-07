"""Local /present run: edit the Allotrope DPPA template IN PLACE for the
Samsung-TTC case, saved to a separate file for comparison with the remote agent's
deck. Focused on the key story slides; preserves masters/logos/diagrams.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parent.parent.parent.parent
# Template is a local-only input (the pristine Allotrope deck); not tracked in git.
TEMPLATE = REPO / "present" / "Allotrope DPPA insights (original template).pptx"
OUTPUT = REPO / "present" / "Allotrope DPPA insights (local run).pptx"
ASSETS = REPO / "present" / "assets_local"  # regenerated each run (gitignored)


def render_charts(out: Path):
    """Render the strike-sweep, adder, and regime PNGs (self-contained, gitignored)."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    out.mkdir(parents=True, exist_ok=True)
    teal, green, red = "#155B55", "#38761D", "#B23A3A"
    plt.rcParams.update({"font.size": 13, "font.family": "DejaVu Sans", "axes.edgecolor": "#888"})

    strikes = ["1012", "1227", "1443", "1658", "1873"]
    buyer = [-25.20, -10.13, 4.95, 20.02, 35.10]
    npv = [-79.7, -74.1, -68.5, -62.9, -57.3]
    fig, ax = plt.subplots(figsize=(7.0, 4.1), dpi=220)
    ax.bar(strikes, buyer, color=[green if v < 0 else red for v in buyer], width=0.6, zorder=3)
    ax.axhline(0, color="#333", lw=1)
    ax.set_ylabel("Buyer vs EVN (B VND/yr)\n(negative = saving)", color=teal)
    ax.set_xlabel("Strike (VND/kWh)")
    ax.tick_params(axis="y", labelcolor=teal)
    ax.grid(axis="y", color="#eee", zorder=0)
    ax2 = ax.twinx()
    ax2.plot(strikes, npv, color=teal, marker="o", lw=2.4, zorder=4)
    ax2.set_ylabel("Developer NPV ($M)", color=teal)
    ax2.tick_params(axis="y", labelcolor=teal)
    ax.set_title("Strike Sweep — Buyer saving vs Developer NPV (directional)", color=teal, fontsize=13, weight="bold")
    fig.tight_layout()
    fig.savefig(out / "strike_sweep.png", facecolor="white")
    plt.close(fig)

    adders = ["0", "262", "523*", "785", "1047"]
    ad = [-61.84, -43.52, -25.20, -6.89, 11.43]
    fig, ax = plt.subplots(figsize=(6.6, 3.9), dpi=220)
    ax.bar(adders, ad, color=[green if v < 0 else red for v in ad], width=0.62, zorder=3)
    ax.axhline(0, color="#333", lw=1)
    ax.set_ylabel("Buyer vs EVN (B VND/yr)", color=teal)
    ax.set_xlabel("DPPA grid-service adder (VND/kWh)   * = inherited base 523")
    ax.grid(axis="y", color="#eee", zorder=0)
    ax.set_title("Dominant lever — DPPA adder flips the buyer near 0.9x (directional)", color=teal, fontsize=12.5, weight="bold")
    fig.tight_layout()
    fig.savefig(out / "adder.png", facecolor="white")
    plt.close(fig)

    regs = ["Decision 963\n(current)", "Decision 14\n(legacy)", "Decree 146\n(2-part trial)"]
    bills = [2036.7, 1993.0, 2401.6]
    fig, ax = plt.subplots(figsize=(6.6, 3.9), dpi=220)
    b = ax.bar(regs, bills, color=[teal, green, red], width=0.6, zorder=3)
    ax.set_ylim(1800, 2500)
    ax.set_ylabel("Samsung EVN bill (B VND/yr)", color=teal)
    ax.grid(axis="y", color="#eee", zorder=0)
    for rect, v, d in zip(b, bills, ["baseline", "-2.1%", "+17.9%"]):
        ax.text(rect.get_x() + rect.get_width() / 2, v + 12, d, ha="center", fontsize=11, color="#333")
    ax.set_title("Regime stress — buyer's EVN outside option (directional)", color=teal, fontsize=12.5, weight="bold")
    fig.tight_layout()
    fig.savefig(out / "regime.png", facecolor="white")
    plt.close(fig)

TEAL = "155B55"
BODY = "222222"
GREEN = "38761D"
RED = "B23A3A"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _bunone(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("buChar", "buAutoNum", "buNone"):
        for el in pPr.findall(f"{{{A}}}{tag}"):
            pPr.remove(el)
    pPr.append(pPr.makeelement(f"{{{A}}}buNone", {}))


def set_text(shape, paras, *, font="Cabin"):
    """paras: list of dicts {t, size, bold, color, align}."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, spec in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        _bunone(p)
        run = p.add_run()
        run.text = spec["t"]
        run.font.name = font
        run.font.size = Pt(spec.get("size", 12))
        run.font.bold = spec.get("bold", False)
        run.font.color.rgb = RGBColor.from_string(spec.get("color", BODY))
        if spec.get("space_after") is not None:
            p.space_after = Pt(spec["space_after"])


def geom(shape):
    return shape.left, shape.top, shape.width, shape.height


def delete_shape(shape):
    shape._element.getparent().remove(shape._element)


def replace_picture(slide, shape, png):
    left, top, width, height = geom(shape)
    delete_shape(shape)
    slide.shapes.add_picture(str(png), left, top, width, height)


def add_textbox(slide, x, y, w, h, paras, font="Cabin"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(tb, paras, font=font)
    return tb


def set_para_text(shape, idx, new_text):
    """Replace one paragraph's text in place, preserving its run formatting
    (color/size/font). Used for template slides so brand styling is untouched."""
    para = shape.text_frame.paragraphs[idx]
    runs = list(para.runs)
    if runs:
        runs[0].text = new_text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        para.add_run().text = new_text


def main():
    render_charts(ASSETS)
    prs = Presentation(str(TEMPLATE))
    s = prs.slides

    # --- Slide 0: title (white text on the dark teal panel) ---
    set_text(s[0].shapes[5], [
        {"t": "Samsung SEVT × TTC Duc Hue 2", "size": 25, "bold": True, "color": "FFFFFF"},
        {"t": "Vietnam's First Grid-Connected DPPA — Directional Economics", "size": 14, "color": "EAF3F1"},
    ])

    # --- Slide 7: strike sweep ---
    set_text(s[7].shapes[4], [{"t": "Strike Sweep: Buyer Saving vs Developer NPV", "size": 22, "bold": True, "color": TEAL}])
    set_text(s[7].shapes[2], [{"t": "At the 1,012 VND/kWh ceiling strike the buyer saves ~25 B VND/yr ($0.95M); the developer is sub-economic across the band — no overlap. Directional.", "size": 11.5, "color": BODY}])
    replace_picture(s[7], s[7].shapes[3], ASSETS / "strike_sweep.png")
    add_textbox(s[7], 7.95, 1.55, 1.95, 3.4, [
        {"t": "AT A GLANCE", "size": 10, "bold": True, "color": TEAL, "space_after": 6},
        {"t": "Buyer eff. cost", "size": 10, "color": "777777"},
        {"t": "1,552 VND/kWh", "size": 14, "bold": True, "color": GREEN, "space_after": 6},
        {"t": "EVN avoided", "size": 10, "color": "777777"},
        {"t": "1,913 VND/kWh", "size": 14, "bold": True, "color": BODY, "space_after": 6},
        {"t": "Developer NPV", "size": 10, "color": "777777"},
        {"t": "−$80M → −$57M", "size": 13, "bold": True, "color": RED, "space_after": 6},
        {"t": "Overlap", "size": 10, "color": "777777"},
        {"t": "none", "size": 14, "bold": True, "color": RED},
    ])

    # --- Slide 8: sensitivities (adder + regime) ---
    # Capture BOTH picture refs before deleting either (deletion shifts indices).
    s8_pic_left = s[8].shapes[4]
    s8_pic_right = s[8].shapes[5]
    set_text(s[8].shapes[0], [{"t": "Sensitivities: DPPA Adder & Tariff Regime", "size": 22, "bold": True, "color": TEAL}])
    set_text(s[8].shapes[3], [{"t": "Two levers move the buyer result: the inherited DPPA grid-service adder (the dominant lever, break-even ~0.9×) and the tariff regime — the Decree 146 two-part tariff (actual bills from Jul 2026) lifts Samsung's EVN bill ~18%. Directional.", "size": 11.5, "color": BODY}])
    replace_picture(s[8], s8_pic_left, ASSETS / "adder.png")
    replace_picture(s[8], s8_pic_right, ASSETS / "regime.png")

    # --- Slide 13: worked case ---
    s13_chart = s[13].shapes[8]  # capture before any deletion
    s13_smallpic = s[13].shapes[7]
    # Shape [9] is the visible top-left title; [2] sits behind the chart.
    set_text(s[13].shapes[9], [{"t": "Samsung × TTC: Worked Case (directional)", "size": 18, "bold": True, "color": TEAL}])
    set_text(s[13].shapes[4], [
        {"t": "THE DEAL", "size": 11, "bold": True, "color": TEAL, "space_after": 3},
        {"t": "Buyer: Samsung SEVT (Thai Nguyen)", "size": 10.5, "color": BODY},
        {"t": "Generator: TTC Duc Hue 2, Tay Ninh", "size": 10.5, "color": BODY},
        {"t": "49 MWp / 41.4 MWac · ~70 GWh/yr", "size": 10.5, "color": BODY},
        {"t": "Financial CfD, Decree 57/2025 · live 1 Jun 2026", "size": 10.5, "color": BODY, "space_after": 8},
        {"t": "DIRECTIONAL RESULT", "size": 11, "bold": True, "color": TEAL, "space_after": 3},
        {"t": "Buyer saves ~25 B VND/yr ($0.95M) at the 1,012 strike", "size": 10.5, "color": GREEN},
        {"t": "Buyer flips to premium above ~1,440 VND/kWh", "size": 10.5, "color": BODY},
        {"t": "Developer sub-economic across band (no overlap)", "size": 10.5, "color": RED},
        {"t": "Adder is the dominant lever (break-even ~0.9×)", "size": 10.5, "color": BODY},
    ])
    replace_picture(s[13], s13_chart, ASSETS / "strike_sweep.png")
    delete_shape(s13_smallpic)
    add_textbox(s[13], 4.14, 4.45, 5.7, 0.95, [
        {"t": "Recommended: buyer-favourable, developer sub-economic — directional, not bankable.", "size": 11, "bold": True, "color": TEAL},
    ])

    # --- Slide 14: recommendation / path to bankable ---
    set_text(s[14].shapes[2], [{"t": "Recommendation & Path to Bankable (directional)", "size": 18, "bold": True, "color": TEAL}])
    set_text(s[14].shapes[4], [
        {"t": "Position: buyer-favourable, developer sub-economic — no clean overlap. Directional, not bankable.", "size": 13, "bold": True, "color": TEAL, "space_after": 8},
        {"t": "For the offtaker (Samsung): the DPPA structurally beats EVN at the regulated ceiling strike (~$0.95M/yr saving on the 70 GWh slice); it strengthens as the Decree 146 two-part tariff hits actual bills from Jul 2026 (Phase 3).", "size": 11.5, "color": BODY, "space_after": 6},
        {"t": "For the generator (TTC): do not read 'uneconomic' from this screen — it is conservative (revenue on the contracted 70 GWh, not the plant's full ~92.5 GWh yield; $750/kW capex). Counting full output + a merchant tail would lift NPV.", "size": 11.5, "color": BODY, "space_after": 6},
        {"t": "Three inputs move this from directional to bankable:", "size": 11.5, "bold": True, "color": BODY, "space_after": 2},
        {"t": "1.  the real DPPA grid / wheeling fee (the dominant lever)", "size": 11, "color": BODY},
        {"t": "2.  a Tay Ninh site-specific CFMP / FMP series", "size": 11, "color": BODY},
        {"t": "3.  the actual negotiated strike and tenor", "size": 11, "color": BODY},
    ])

    # --- Currency fixes to retained template slides (DPPA news, mid-2026) ---
    # Slide 1: a live first implementation now exists (contradicted "still emerging").
    set_para_text(s[1].shapes[17], 0, "First grid-connected DPPA now live (Samsung–TTC, Jun 2026); broader rollout still early")
    set_para_text(s[1].shapes[19], 0, "First tangible project example now operating, reducing cost-saving uncertainty")
    # Slide 4: grid-connected CfD settles vs the spot market price (FMP), not retail
    # (matches slide 11; the live mechanism under Decree 57).
    set_para_text(s[4].shapes[5], 2, "On a monthly basis, the energy customer settles the financial payment with the RE developer equal to the difference between the spot market price (FMP) and the strike price")
    # Slide 6: a first-in-market grid example now exists.
    set_para_text(s[6].shapes[6], 0, "A first grid-connected example is now live (Samsung–TTC, Jun 2026); the market is still thin, so early buyers remain first movers")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"saved: {OUTPUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
