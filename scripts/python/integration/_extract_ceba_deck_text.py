"""Extract text from a deck (.pptx) into a flat text file for review.

Decks:
  - CEBA DPPA 2026.pptx                   → ceba_dppa_2026_text.txt          (default)
  - DPPA July 2026 Case Studies.pptx     → dppa_july_2026_case_studies_text.txt

Usage (from repo root):
    .venv/Scripts/python.exe scripts/python/integration/_extract_ceba_deck_text.py
    .venv/Scripts/python.exe scripts/python/integration/_extract_ceba_deck_text.py --deck july
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parents[3]


def _walk_shapes(shapes, depth: int = 0) -> list:
    """Recursively collect all shapes (including grouped)."""
    out = []
    for shape in shapes:
        out.append((shape, depth))
        if shape.shape_type == 6:  # GROUP
            out.extend(_walk_shapes(shape.shapes, depth + 1))
    return out


def extract_text(pptx_path: Path) -> str:
    prs = Presentation(str(pptx_path))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"[Slide {i}]")
        for shape, _ in _walk_shapes(slide.shapes):
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        lines.append(text)
            if shape.has_table:
                tbl = shape.table
                for row in tbl.rows:
                    cells = []
                    for cell in row.cells:
                        cell_text = " ".join(
                            p.text for p in cell.text_frame.paragraphs if p.text.strip()
                        )
                        cells.append(cell_text.strip())
                    lines.append(" | ".join(cells))
        lines.append("")
    return "\n".join(lines)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--deck",
        choices=("ceba", "july"),
        default="ceba",
        help="Which deck to extract (default: ceba).",
    )
    parser.add_argument(
        "--src",
        type=Path,
        default=None,
        help="Override the source deck path (default: ceba-review/<deck>.pptx).",
    )
    parser.add_argument(
        "--out",
        type=Path,
        default=None,
        help="Override the output text path.",
    )
    args = parser.parse_args(argv)

    # Local import so this script is runnable without configuring sys.path.
    SCRIPTS_PYTHON = REPO_ROOT / "scripts" / "python"
    if str(SCRIPTS_PYTHON) not in sys.path:
        sys.path.insert(0, str(SCRIPTS_PYTHON))
    from integration.ceba_deck.deck_config import get_deck  # noqa: E402

    config = get_deck(args.deck)
    src = args.src or config.source_pptx
    out = args.out or config.text_txt

    if not src.exists():
        print(f"Deck not found: {src}", file=sys.stderr)
        return 1
    text = extract_text(src)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(text, encoding="utf-8")
    print(
        f"Extracted {len(text):,} chars from {src.name} → {out.relative_to(REPO_ROOT)}"
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
