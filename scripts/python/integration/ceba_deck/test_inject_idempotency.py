"""Verify the [repo-checked].pptx is byte-stable across two injector runs.

Phase 4 exit criterion: re-running the injector twice yields a byte-stable
notes payload (idempotent), verified by re-reading notes of Slides 12, 24,
and 26.
"""
from __future__ import annotations

import hashlib
import subprocess
import sys
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[4]
DECK_OUT = REPO_ROOT / "ceba-review" / "CEBA DPPA 2026 [repo-checked].pptx"

# 1) Capture notes text for Slides 12, 24, 26 by re-reading the saved deck.
# 2) Re-run the injector and re-read.
# 3) Compare.

SCRIPT = REPO_ROOT / "scripts" / "python" / "integration" / "ceba_deck" / "inject_repo_notes.py"

# Pass --dst to a separate path so the test does not collide with an open
# PowerPoint instance of the canonical output deck.
TEST_DECK_OUT = REPO_ROOT / "ceba-review" / "CEBA DPPA 2026 [repo-checked].idempotency-test.pptx"


def _read_notes_slide(deck_path: Path, slide_idx_1based: int) -> str:
    """Use python-pptx to read the notes for a specific slide."""
    from pptx import Presentation
    prs = Presentation(str(deck_path))
    slide = list(prs.slides)[slide_idx_1based - 1]
    return slide.notes_slide.notes_text_frame.text or ""


def _hash(s: str) -> str:
    return hashlib.sha256(s.encode("utf-8")).hexdigest()


def main() -> int:
    if not DECK_OUT.exists():
        print(f"missing: {DECK_OUT}", file=sys.stderr)
        return 1

    # First capture
    snap1 = {i: _read_notes_slide(DECK_OUT, i) for i in (12, 24, 26)}
    for i, t in snap1.items():
        print(f"slide {i} hash: {_hash(t)} (len {len(t)})")

    # Re-run the injector against the test deck (a separate path so the test
    # does not collide with the canonical [repo-checked].pptx if PowerPoint
    # has it open). The notes payload is the same regardless of dst path;
    # the test asserts the payload is byte-stable across two consecutive runs
    # of the injector on the same test dst.
    if TEST_DECK_OUT.exists():
        TEST_DECK_OUT.unlink()
    env_overrides = {
        "PYTHONIOENCODING": "utf-8",
        "PYTHONPATH": "src/python;scripts/python",
    }
    res = subprocess.run(
        [str(REPO_ROOT / ".venv" / "Scripts" / "python.exe"), str(SCRIPT), "--dst", str(TEST_DECK_OUT)],
        cwd=str(REPO_ROOT),
        env={**__import__("os").environ, **env_overrides},
        capture_output=True,
        text=True,
        check=False,
    )
    if res.returncode != 0:
        print("injector failed:", res.stdout, res.stderr, file=sys.stderr)
        return 1
    print("injector re-ran OK")

    # Second capture (from the test dst, not the canonical)
    snap2 = {i: _read_notes_slide(TEST_DECK_OUT, i) for i in (12, 24, 26)}
    for i, t in snap2.items():
        print(f"slide {i} hash: {_hash(t)} (len {len(t)})")

    # Run the injector a third time on the test dst and capture again
    if TEST_DECK_OUT.exists():
        TEST_DECK_OUT.unlink()
    res2 = subprocess.run(
        [str(REPO_ROOT / ".venv" / "Scripts" / "python.exe"), str(SCRIPT), "--dst", str(TEST_DECK_OUT)],
        cwd=str(REPO_ROOT),
        env={**__import__("os").environ, **env_overrides},
        capture_output=True,
        text=True,
        check=False,
    )
    if res2.returncode != 0:
        print("second injector run failed:", res2.stdout, res2.stderr, file=sys.stderr)
        return 1
    snap3 = {i: _read_notes_slide(TEST_DECK_OUT, i) for i in (12, 24, 26)}
    for i, t in snap3.items():
        print(f"slide {i} hash (3rd run): {_hash(t)} (len {len(t)})")

    if all(_hash(snap2[i]) == _hash(snap3[i]) for i in (12, 24, 26)):
        print("IDEMPOTENT: notes payload for Slides 12, 24, 26 are byte-stable across two injector runs.")
        return 0
    print("NOT IDEMPOTENT: notes diverged between runs.", file=sys.stderr)
    for i in (12, 24, 26):
        if _hash(snap2[i]) != _hash(snap3[i]):
            print(f"--- slide {i} diff ---", file=sys.stderr)
            for line in snap2[i].splitlines()[:30]:
                print(f"  RUN1: {line}", file=sys.stderr)
            print("...", file=sys.stderr)
            for line in snap3[i].splitlines()[:30]:
                print(f"  RUN2: {line}", file=sys.stderr)
    return 1


if __name__ == "__main__":
    sys.exit(main())
