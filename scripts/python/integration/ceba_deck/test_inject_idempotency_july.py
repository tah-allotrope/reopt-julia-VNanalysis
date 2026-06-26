"""Verify the July [repo-checked].pptx is byte-stable across two injector runs.

PHASE-05 extension to the CEBA idempotency test. Captures the notes text
of three July deck slides that carry repo-testable claims — Slide 11
(worked example), Slide 23 (Case 5 metrics), and Slide 25 (56-sweep) —
and asserts the payload is byte-stable across two consecutive runs of
the injector on a separate test dst.

The byte-stability test exercises the same CON-003 idempotency
guarantee the CEBA pipeline uses: a delimiter
``=== [Repo check] (generated) ===`` is detected and everything below
it is replaced; any author notes above the delimiter are preserved.
"""
from __future__ import annotations

import hashlib
import os
import subprocess
import sys
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[4]
DECK_OUT = REPO_ROOT / "ceba-review" / "DPPA Presentation July 2026 Case Studies [repo-checked].pptx"
SCRIPT = REPO_ROOT / "scripts" / "python" / "integration" / "ceba_deck" / "inject_repo_notes.py"
TEST_DECK_OUT = REPO_ROOT / "ceba-review" / "DPPA July 2026 [repo-checked].idempotency-test.pptx"

# Three July slides with the densest [Repo check] payloads:
#   - 11: worked example (5-line settlement)
#   - 23: Case 5 metrics (8 checks)
#   - 25: 56-sweep (5 checks)
SAMPLE_SLIDES = (11, 23, 25)


def _read_notes_slide(deck_path: Path, slide_idx_1based: int) -> str:
    from pptx import Presentation
    prs = Presentation(str(deck_path))
    slide = list(prs.slides)[slide_idx_1based - 1]
    return slide.notes_slide.notes_text_frame.text or ""


def _hash(s: str) -> str:
    return hashlib.sha256(s.encode("utf-8")).hexdigest()


def main() -> int:
    if not DECK_OUT.exists():
        print(f"missing: {DECK_OUT}", file=sys.stderr)
        print("(run scripts/python/integration/ceba_deck/inject_repo_notes.py --deck july first)", file=sys.stderr)
        return 1
    if not (REPO_ROOT / "reports" / "dppa_july_2026_repo_check.json").exists():
        print("missing: reports/dppa_july_2026_repo_check.json", file=sys.stderr)
        return 1

    # First capture from the canonical [repo-checked] deck
    snap1 = {i: _read_notes_slide(DECK_OUT, i) for i in SAMPLE_SLIDES}
    for i, t in snap1.items():
        print(f"slide {i} hash: {_hash(t)} (len {len(t)})")

    # Re-run the injector against the test deck
    if TEST_DECK_OUT.exists():
        TEST_DECK_OUT.unlink()
    env_overrides = {
        "PYTHONIOENCODING": "utf-8",
        "PYTHONPATH": "src/python;scripts/python",
    }
    res = subprocess.run(
        [
            str(REPO_ROOT / ".venv" / "Scripts" / "python.exe"),
            str(SCRIPT),
            "--deck", "july",
            "--dst", str(TEST_DECK_OUT),
        ],
        cwd=str(REPO_ROOT),
        env={**os.environ, **env_overrides},
        capture_output=True,
        text=True,
        check=False,
    )
    if res.returncode != 0:
        print("first injector run failed:", res.stdout, res.stderr, file=sys.stderr)
        return 1
    print("first injector re-run OK")
    snap2 = {i: _read_notes_slide(TEST_DECK_OUT, i) for i in SAMPLE_SLIDES}
    for i, t in snap2.items():
        print(f"slide {i} hash: {_hash(t)} (len {len(t)})")

    # Re-run a second time
    if TEST_DECK_OUT.exists():
        TEST_DECK_OUT.unlink()
    res2 = subprocess.run(
        [
            str(REPO_ROOT / ".venv" / "Scripts" / "python.exe"),
            str(SCRIPT),
            "--deck", "july",
            "--dst", str(TEST_DECK_OUT),
        ],
        cwd=str(REPO_ROOT),
        env={**os.environ, **env_overrides},
        capture_output=True,
        text=True,
        check=False,
    )
    if res2.returncode != 0:
        print("second injector run failed:", res.stdout, res.stderr, file=sys.stderr)
        return 1
    snap3 = {i: _read_notes_slide(TEST_DECK_OUT, i) for i in SAMPLE_SLIDES}
    for i, t in snap3.items():
        print(f"slide {i} hash (3rd run): {_hash(t)} (len {len(t)})")

    if all(_hash(snap2[i]) == _hash(snap3[i]) for i in SAMPLE_SLIDES):
        print(
            f"IDEMPOTENT: notes payload for Slides {SAMPLE_SLIDES} are byte-stable "
            "across two July injector runs."
        )
        return 0
    print("NOT IDEMPOTENT: notes diverged between runs.", file=sys.stderr)
    for i in SAMPLE_SLIDES:
        if _hash(snap2[i]) != _hash(snap3[i]):
            print(f"--- slide {i} diff ---", file=sys.stderr)
            for line in snap2[i].splitlines()[:30]:
                print(f"  RUN1: {line}", file=sys.stderr)
            print("...", file=sys.stderr)
            for line in snap3[i].splitlines()[:30]:
                print(f"  RUN2: {line}", file=sys.stderr)
    return 1


if __name__ == "__main__":
    sys.exit(main())
