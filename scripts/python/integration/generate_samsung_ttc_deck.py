"""Generate the Samsung-TTC DPPA presentation IN PLACE on the Allotrope template.

The named template `present/Allotrope DPPA insights.pptx` is not present in this
checkout, so we use the canonical Allotrope-branded template
(`reports/decks/conformance/template/allotrope-template.pptx`) as the base — the
same master/layout/logo set the production decks are built from
(see `reports/decks/2026-04-21-dppa-case-3-allotrope.pptx`).

Approach: preserve masters, layouts, logos, banner, background, fonts, and colors.
Only the editable text frames, table cells, and chart content are swapped to the
Samsung SEVT -- TTC Duc Hue 2 DPPA story. Every headline number is labeled
DIRECTIONAL (undisclosed strike + proxy CFMP market reference).

Data sources:
  artifacts/reports/samsung_ttc/*.json (plan PHASE 01-04 outputs)
  research/2026-06-04_samsung-ttc-dppa.md (deal + regulatory context)
"""
import copy
import json
import os

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

REPO = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))
TEMPLATE = os.path.join(REPO, "reports", "decks", "conformance", "template", "allotrope-template.pptx")
ART = os.path.join(REPO, "artifacts", "reports", "samsung_ttc")
OUT_DIR = os.path.join(REPO, "present")
OUT = os.path.join(OUT_DIR, "Allotrope DPPA insights.pptx")

PRIMARY_TEAL = RGBColor(0x15, 0x5B, 0x55)
ACCENT_GREEN = RGBColor(0x38, 0x76, 0x1D)
SAVE_TEAL = RGBColor(0x14, 0x76, 0x68)
PREMIUM_RED = RGBColor(0xC0, 0x39, 0x2B)


def load(name):
    with open(os.path.join(ART, name), "r", encoding="utf-8") as fh:
        return json.load(fh)


# ---------------------------------------------------------------- helpers
def shape_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def _set_run_text(para, text):
    """Set paragraph text on its first run, preserving that run's formatting."""
    runs = para.runs
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        run = para.add_run()
        run.text = text


def set_tf(tf, texts):
    """Map a list of paragraph strings onto a text frame, preserving formatting.

    Reuses existing paragraphs (keeps bullet/indent/font). Extra paragraphs are
    cloned from the last template paragraph; surplus template paragraphs blanked.
    """
    if isinstance(texts, str):
        texts = [texts]
    paras = tf.paragraphs
    n = len(paras)
    for i, txt in enumerate(texts):
        if i < n:
            _set_run_text(paras[i], txt)
        else:
            template_p = paras[n - 1]._p
            new_p = copy.deepcopy(template_p)
            template_p.addnext(new_p)
            paras = tf.paragraphs  # refresh
            _set_run_text(paras[i], txt)
    # blank any leftover template paragraphs
    for j in range(len(texts), len(tf.paragraphs)):
        _set_run_text(tf.paragraphs[j], "")


def set_cell(cell, text, bold=None, color=None):
    if cell.is_spanned:
        return
    tf = cell.text_frame
    p = tf.paragraphs[0]
    _set_run_text(p, text)
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    if p.runs:
        run = p.runs[0]
        if bold is not None:
            run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color


def style_chart(chart, font_pt=9):
    chart.has_title = False
    try:
        chart.font.size = Pt(font_pt)
        chart.font.name = "Cabin"
    except (AttributeError, TypeError):
        # Some chart types expose no font attribute; leave them unstyled.
        pass


# ---------------------------------------------------------------- load data
defn = load("2026-06-04_samsung-ttc_dppa-definition.json")
combined = load("2026-06-04_samsung-ttc_combined-decision.json")
slice_ = load("2026-06-04_samsung-ttc_contracted-slice.json")["contracted_slice"]
solar = load("2026-06-04_samsung-ttc_solar-summary.json")
strike = load("2026-06-04_samsung-ttc_strike-sensitivity.json")
contract_risk = load("2026-06-04_samsung-ttc_contract-risk.json")
regime = load("2026-06-04_samsung-ttc_regime-stress.json")

B = 1e9  # VND -> billion VND
sweep = strike["sweep"]
adder = contract_risk["adder_sensitivity"]["results"]
regimes = regime["regimes"]

prs = Presentation(TEMPLATE)
slides = list(prs.slides)

# ================================================================ SLIDE 1 — title
s = slides[0]
set_tf(shape_by_name(s, "Google Shape;272;p52").text_frame, [
    "Allotrope DPPA Insights",
    "Samsung SEVT — TTC Duc Hue 2: Vietnam’s First Grid-Connected DPPA",
    "June 2026  —  Directional economics (illustrative, not bankable)",
])

# ================================================================ SLIDE 2 — deal at a glance (8x3 merged table)
s = slides[1]
set_tf(shape_by_name(s, "Google Shape;278;p53").text_frame, [
    "Vietnam’s First Grid-Connected DPPA: ",
    "Deal at a Glance  (live 1 Jun 2026)",
])
tbl = None
for sh in s.shapes:
    if sh.has_table:
        tbl = sh.table
hdr = ["Group", "Parameter", "Value / Basis"]
for j, h in enumerate(hdr):
    set_cell(tbl.cell(0, j), h, bold=True)
rows = [
    # (col0-group-origin or None, col1 param, col2 value)
    ("Parties", "Buyer", "Samsung Electronics Vietnam Thai Nguyen (SEVT) — Thai Nguyen, north"),
    (None,      "Generator", "TTC Duc Hue 2 solar — Tay Ninh, south (TTC Duc Hue–Long An Power JSC)"),
    ("Plant",   "Capacity", "49 MWp DC / 41.4 MWac — ground-mount, no storage"),
    (None,      "Output", "~70 GWh/yr contracted (~19.3% AC capacity factor); ~46,000 t CO₂ avoided/yr"),
    (None,      "Milestones", "COD 19 May 2026; DPPA live 1 Jun 2026"),
    ("Contract","Mechanism", "Grid-connected DPPA under Decree 57/2025 (three-contract model)"),
    (None,      "Settlement", "Financial contract-for-differences (CfD) vs market reference"),
]
for i, (grp, param, val) in enumerate(rows, start=1):
    if grp is not None:
        set_cell(tbl.cell(i, 0), grp, bold=True, color=PRIMARY_TEAL)
    set_cell(tbl.cell(i, 1), param, bold=True)
    set_cell(tbl.cell(i, 2), val)

# ================================================================ SLIDE 3 — deal information (11x2 table)
s = slides[2]
set_tf(shape_by_name(s, "Google Shape;289;p54").text_frame, ["DEAL INFORMATION & STRIKE BASIS"])
# remove right-side content pictures (keep top-right logo 288)
for nm in ("Google Shape;285;p54", "Google Shape;287;p54", "Google Shape;286;p54"):
    sh = shape_by_name(s, nm)
    if sh is not None:
        remove_shape(sh)
tbl = None
for sh in s.shapes:
    if sh.has_table:
        tbl = sh.table
info = [
    ("Buyer", "Samsung Electronics Vietnam Thai Nguyen (SEVT)"),
    ("Buyer location", "Yen Binh Industrial Park, Thai Nguyen province (north)"),
    ("Generator", "TTC Duc Hue 2 solar farm, Tay Ninh province (south)"),
    ("Developer", "TTC Duc Hue–Long An Power JSC"),
    ("Capacity", "49 MWp DC / 41.4 MWac, ground-mounted, no storage"),
    ("Annual volume", "~70 GWh/yr to buyer (~7% of synthetic SEVT load)"),
    ("Mechanism", "Grid-connected DPPA, national grid, Decree 57/2025"),
    ("Settlement", "Financial CfD; generator sells to VWEM spot, settles difference vs SEVT"),
    ("Strike basis", "Anchored to Southern ground-mount ceiling ≈ 1,012 VND/kWh (Decree 57 cap)"),
    ("Market ref", "Proxy CFMP/FMP series (transferred, not site-specific)"),
    ("Disclosure", "Strike, tenor, KPP/grid fee UNDISCLOSED — all numbers DIRECTIONAL"),
]
for i, (k, v) in enumerate(info):
    set_cell(tbl.cell(i, 0), k, bold=True, color=PRIMARY_TEAL)
    set_cell(tbl.cell(i, 1), v)

# ================================================================ SLIDE 4 — key findings (5 cards)
s = slides[3]
card_titles = [
    ("Google Shape;310;p55", "Buyer Saves (Directional)"),
    ("Google Shape;311;p55", "Developer Sub-Economic"),
    ("Google Shape;312;p55", "Adder Is the Dominant Lever"),
    ("Google Shape;313;p55", "Regime Stress"),
    ("Google Shape;314;p55", "Recommendation"),
]
for nm, t in card_titles:
    set_tf(shape_by_name(s, nm).text_frame, [t])
buyer_save_b = -sweep[0]["buyer_minus_benchmark_vnd"] / B  # 25.20
npv_lo = sweep[0]["developer_npv_usd"] / 1e6
npv_hi = sweep[-1]["developer_npv_usd"] / 1e6
dec146 = next(r for r in regimes if "146" in r["regime_id"])
dec14 = next(r for r in regimes if r["regime_id"] == "decision_14_2025_legacy")
card_bodies = [
    ("Google Shape;315;p55",
     "At the Southern-ceiling strike (1,012 VND/kWh) the buyer saves ~{:.1f} B VND/yr "
     "(~${:.2f}M) on the 70 GWh slice: effective {:,.0f} vs EVN avoided {:,.0f} VND/kWh. Directional.".format(
         buyer_save_b, slice_["buyer_savings_usd"] / 1e6,
         slice_["buyer_effective_cost_vnd_per_kwh"], slice_["evn_avoided_cost_vnd_per_kwh"])),
    ("Google Shape;316;p55",
     f"PySAM Single Owner NPV runs −${abs(npv_lo):.1f}M → −${abs(npv_hi):.1f}M across the strike band (IRR null) on a "
     "conservative 70 GWh / $750/kW basis. No buyer–developer overlap."),
    ("Google Shape;317;p55",
     "The DPPA grid-service adder (523 VND/kWh, inherited) is the single biggest lever on buyer cost; "
     "the buyer flips from saving to a premium near ~0.9× that adder."),
    ("Google Shape;318;p55",
     "Decree 146 two-part tariff trial lifts the buyer’s EVN bill +{:.1f}% (+{:.0f} B VND); "
     "Decision 14 legacy is {:.1f}%. A higher bill raises DPPA value unless it double-charges DPPA volume.".format(
         dec146["delta_pct"], dec146["annual_bill_delta_vnd"] / B, dec14["delta_pct"])),
    ("Google Shape;319;p55",
     "Directional position: buyer-favorable, developer sub-economic. Every figure rides on an undisclosed "
     "strike and a proxy CFMP series — directional, not bankable."),
]
for nm, body in card_bodies:
    set_tf(shape_by_name(s, nm).text_frame, [body])

# ================================================================ SLIDE 5 — plant & solar profile (6x2 table)
s = slides[4]
set_tf(shape_by_name(s, "Google Shape;332;p56").text_frame, [
    "TTC Duc Hue 2: ",
    "Plant & Solar Profile (directional)",
])
# remove the stale Ninh Sim site image; replace with a clean directional note chart later not needed
img = shape_by_name(s, "Google Shape;331;p56")
if img is not None:
    remove_shape(img)
tbl = None
for sh in s.shapes:
    if sh.has_table:
        tbl = sh.table
plant = [
    ("Plant type", "49 MWp DC / 41.4 MWac ground-mounted solar, no storage"),
    ("Location", "Tay Ninh province, southern Vietnam (high-irradiance)"),
    ("Annual output", "~70 GWh/yr contracted to SEVT (CfD volume)"),
    ("Capacity factor", "~{:.1f}% AC; peak ~{:.1f} MWac (no clipping)".format(
        solar["ac_capacity_factor"] * 100.0, solar["peak_ac_kw"] / 1000.0)),
    ("Solar profile", "PySAM PVWatts v8 on cached southern Himawari resource, scaled to 70 GWh"),
    ("CO₂ avoided", "~46,000 tonnes per year"),
]
for i, (k, v) in enumerate(plant):
    set_cell(tbl.cell(i, 0), k, bold=True, color=PRIMARY_TEAL)
    set_cell(tbl.cell(i, 1), v)

# add an illustrative monthly-shape-free KPI panel via a small native chart (annual energy split)
cd = CategoryChartData()
cd.categories = ["Contracted to SEVT", "Plant native yield"]
cd.add_series("GWh/yr", (70.0, round(solar["native_annual_gwh"], 1)))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.3), Inches(1.25),
                        Inches(3.45), Inches(3.9), cd)
style_chart(gf.chart)
gf.chart.has_legend = False
plot = gf.chart.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0'
plot.data_labels.number_format_is_linked = False

# ================================================================ SLIDE 6 — strike sweep chart
s = slides[5]
set_tf(shape_by_name(s, "Google Shape;337;p57").text_frame, [
    "Strike Sweep: Buyer Premium & Developer NPV",
])
set_tf(shape_by_name(s, "Google Shape;338;p57").text_frame, [
    "Directional: the buyer saves vs EVN below ~1,440 VND/kWh; the developer NPV stays negative across the "
    "whole Southern-ceiling → EVN-avoided band (no overlap). Strike basis + proxy CFMP labeled.",
])
# remove stale financial chart pictures + group
for nm in ("Google Shape;352;p57", "Google Shape;353;p57", "Google Shape;355;p57", "Google Shape;339;p57"):
    sh = shape_by_name(s, nm)
    if sh is not None:
        remove_shape(sh)
cats = ["{:,.0f}".format(p["strike_vnd_per_kwh"]) for p in sweep]
# buyer net cost vs EVN (B VND): negative = buyer saves
buyer_net = [round(p["buyer_minus_benchmark_vnd"] / B, 2) for p in sweep]
cd = CategoryChartData()
cd.categories = cats
cd.add_series("Buyer net vs EVN (B VND/yr; − = saves)", buyer_net)
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.36), Inches(1.95),
                        Inches(4.7), Inches(3.2), cd)
style_chart(gf.chart)
gf.chart.has_legend = True
gf.chart.legend.position = XL_LEGEND_POSITION.BOTTOM
gf.chart.legend.include_in_layout = False

dev_npv = [round(p["developer_npv_usd"] / 1e6, 1) for p in sweep]
cd2 = CategoryChartData()
cd2.categories = cats
cd2.add_series("Developer NPV ($M)", dev_npv)
gf2 = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(5.15), Inches(1.95),
                         Inches(4.55), Inches(3.2), cd2)
style_chart(gf2.chart)
gf2.chart.has_legend = True
gf2.chart.legend.position = XL_LEGEND_POSITION.BOTTOM
gf2.chart.legend.include_in_layout = False
# axis caption via existing footer remains

# ================================================================ SLIDE 7 — methodology
s = slides[6]
set_tf(shape_by_name(s, "Google Shape;360;p58").text_frame, [
    "Methodology: ",
    "Four-Phase Directional DPPA Screen",
])
ph = shape_by_name(s, "Google Shape;363;p58")
if ph is not None:
    remove_shape(ph)
body = shape_by_name(s, "Google Shape;362;p58")
# widen body to full slide since photo removed
body.left = Inches(0.3)
body.width = Inches(9.4)
set_tf(body.text_frame, [
    "We reuse the repo’s tested synthetic/financial DPPA settlement engine (ninhsim Case 2) with the "
    "plant pinned to the disclosed 49 MWp / 70 GWh and no storage — a fixed-plant screen, not a re-optimization.",
    "Phase 1 — Case definition: disclosed deal facts, synthetic SEVT load (~1,000 GWh; 70 GWh ≈ 7% RE), "
    "strike anchored to the Southern ground-mount ceiling (1,012 VND/kWh).",
    "Phase 2 — Solar + settlement: 49 MWp 8760 via PySAM PVWatts v8 (southern resource, scaled to 70 GWh); "
    "hourly CfD settlement and EVN avoided-cost benchmark on the matched slice.",
    "Phase 3 — Strike sweep + developer screen: buyer premium across the strike band and PySAM Single Owner "
    "after-tax IRR/NPV at the fixed 49 MWp; plus a DPPA grid-service adder sensitivity.",
    "Phase 4 — Regime stress + decision: GAP-05 toggle (Decision 963 vs Decision 14 vs Decree 146 two-part) "
    "and a combined, explicitly caveated recommendation.",
    "Every headline number is flagged DIRECTIONAL with explicit strike + market-reference (proxy CFMP) basis.",
])

# ================================================================ SLIDE 8 — adder sensitivity chart
s = slides[7]
set_tf(shape_by_name(s, "Google Shape;368;p59").text_frame, [
    "The Dominant Lever: ",
    "DPPA Grid-Service Adder Sensitivity",
])
ph = shape_by_name(s, "Google Shape;370;p59")
if ph is not None:
    remove_shape(ph)
txt = shape_by_name(s, "Google Shape;369;p59")
txt.left = Inches(0.3)
txt.width = Inches(9.4)
break_even = "~0.9×"
set_tf(txt.text_frame, [
    "The DPPA grid-service adder — inherited from the Case-2 default at 523.34 VND/kWh — is the single "
    "biggest swing factor on the buyer’s effective cost, larger than the strike itself within the regulated band.",
    "Sweeping the adder from 0 to 2×: the buyer’s annual position vs EVN moves from "
    "−{:.1f} B VND (saves) at 0× to +{:.1f} B VND (premium) at 2×. Break-even is near {} the inherited adder.".format(
        -adder[0]["buyer_minus_benchmark_vnd"] / B, adder[-1]["buyer_minus_benchmark_vnd"] / B, break_even),
    "Implication: the adder / KPP grid fee must be deal-calibrated before the sign of the buyer result can be trusted. "
    "Directional only.",
])
acats = ["{:.1f}×\n({:,.0f})".format(r["adder_multiplier"], r["dppa_adder_vnd_per_kwh"]) for r in adder]
avals = [round(r["buyer_minus_benchmark_vnd"] / B, 2) for r in adder]
cd = CategoryChartData()
cd.categories = acats
cd.add_series("Buyer net vs EVN (B VND/yr; − = saves)", avals)
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(2.55),
                        Inches(8.8), Inches(2.6), cd)
style_chart(gf.chart)
gf.chart.has_legend = True
gf.chart.legend.position = XL_LEGEND_POSITION.BOTTOM
gf.chart.legend.include_in_layout = False
plot = gf.chart.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0'
plot.data_labels.number_format_is_linked = False

# ================================================================ SLIDE 9 — regime stress chart
s = slides[8]
set_tf(shape_by_name(s, "Google Shape;377;p60").text_frame, [
    "Regime Stress: ",
    "Buyer EVN Bill Across Tariff Regimes",
])
ph = shape_by_name(s, "Google Shape;376;p60")
if ph is not None:
    remove_shape(ph)
set_tf(shape_by_name(s, "Google Shape;375;p60").text_frame, [
    "The buyer’s EVN bill is the DPPA’s outside option: a higher bill raises the avoided-cost value of "
    "the contract. We stress the synthetic SEVT load across three regimes via the GAP-05 toggle.",
    "Decision 963/2026 (current): baseline. Decision 14/2025 legacy: {:.1f}%. "
    "Decree 146/2025 two-part trial: +{:.1f}% (+{:.0f} B VND) from the capacity charge.".format(
        dec14["delta_pct"], dec146["delta_pct"], dec146["annual_bill_delta_vnd"] / B),
])
set_tf(shape_by_name(s, "Google Shape;378;p60").text_frame, [
    "Caveat: ",
    "the Decree 146 two-part trial is a paper trial, not yet on actual bills, and it risks double-charging "
    "capacity on the DPPA volume — which would erode the buyer savings shown elsewhere in this deck. Directional.",
])
rcats = ["Decision 963\n(current)", "Decision 14\n(legacy)", "Decree 146\n(two-part trial)"]
rvals = [round(r["annual_bill_gvnd"], 1) for r in regimes]
cd = CategoryChartData()
cd.categories = rcats
cd.add_series("Annual EVN bill (B VND/yr)", rvals)
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(5.05), Inches(1.25),
                        Inches(4.6), Inches(3.9), cd)
style_chart(gf.chart)
gf.chart.has_legend = False
plot = gf.chart.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0'
plot.data_labels.number_format_is_linked = False

# ================================================================ SLIDE 10 — Decree 57 / market context
s = slides[9]
set_tf(shape_by_name(s, "Google Shape;383;p61").text_frame, [
    "Context: ",
    "Decree 57/2025 & the Vietnam DPPA Market",
])
set_tf(shape_by_name(s, "Google Shape;385;p61").text_frame, [
    "Decree 57/2025 grid-connected DPPA — a three-contract model:",
    "1) Forward CfD between generator and buyer; 2) generator sells physical output to the VWEM spot market; "
    "3) buyer continues to buy from EVN and settles the difference. The two sites are ~1,500 km apart "
    "(buyer north, plant south), the definitional signature of a financial (not private-wire) DPPA.",
    "Forward price is left to bilateral negotiation but capped at the RE-type ceiling generation tariff — "
    "Southern ground-mount, no storage ≈ 1,012 VND/kWh — which anchors our strike floor.",
    "Buyer avoided cost: EVN production tariff (base 2,204 VND/kWh × ~0.85 standard-hour multiplier "
    "≈ 1,873–1,913 VND/kWh), since solar delivers in standard TOU hours under Decision 963.",
    "Forward risk: the Decree 146 two-part (capacity + energy) tariff trial can double-charge capacity on DPPA "
    "volume and erode buyer savings — the key regime to watch.",
])

# ================================================================ SLIDE 11 — risks & open questions
s = slides[10]
set_tf(shape_by_name(s, "Google Shape;390;p62").text_frame, [
    "Risks & Open Questions: ",
    "What Makes These Numbers Directional",
])
set_tf(shape_by_name(s, "Google Shape;393;p62").text_frame, [
    "Why directional, not bankable:",
    "Strike price is undisclosed — anchored to the Southern ceiling (1,012 VND/kWh) and swept to EVN avoided "
    "cost (~1,873). The buyer result’s sign depends on where the real strike sits.",
    "Market reference (CfD leg) uses a proxy CFMP/FMP series transferred from another site, not a Duc Hue 2 series.",
    "The DPPA grid-service adder / KPP grid fee is inherited from Case 2 (523 VND/kWh) and is the dominant lever — "
    "it must be deal-calibrated.",
    "Solar 8760 is from a representative southern resource scaled to 70 GWh, not a site-specific Tay Ninh dataset; "
    "developer NPV uses conservative $750/kW capex on the contracted 70 GWh (excludes merchant tail).",
])
set_tf(shape_by_name(s, "Google Shape;392;p62").text_frame, [
    "Recommended position (directional):",
    "Buyer-favorable, developer sub-economic on the contracted slice — no buyer–developer overlap under "
    "conservative assumptions. To move toward bankable: source the actual strike & tenor, a Duc Hue 2 FMP/CFMP "
    "series, the negotiated KPP/grid fee, and SEVT’s real 8760 load and connection voltage.",
])

# ================================================================ SLIDE 12 — closing / sources
s = slides[11]
set_tf(shape_by_name(s, "Google Shape;402;p63").text_frame, [
    "Allotrope — DPPA Insights",
    "Sources: pv-magazine, TheInvestor.vn, VnEconomy, Reccessary;",
    "Norton Rose Fulbright & A&O Shearman (Decree 57); EVN/MOIT solar ceilings.",
    "All deal economics DIRECTIONAL: undisclosed strike + proxy CFMP basis.",
])

# ---------------------------------------------------------------- save
os.makedirs(OUT_DIR, exist_ok=True)
prs.save(OUT)
print(f"Saved {OUT} ({len(list(prs.slides))} slides)")
